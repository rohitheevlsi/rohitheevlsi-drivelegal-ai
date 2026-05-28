# generate_presentation_7.py
"""Generate a clean, 7‑slide PowerPoint deck for the DriveLegal AI hackathon.

Features applied:
- Light‑gray background for all slides (consistent theme)
- Arial / Helvetica fonts, dark‑gray text for readability
- Minimal but meaningful image on the demo slide
- Simple bullet styling and centered titles
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# -------------------------------------------------------------------
# Theme helpers
# -------------------------------------------------------------------
def apply_theme(slide, title_text=None):
    """Apply a light‑gray background and optional title styling.
    """
    # Background colour
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)  # #F5F5F5
    # Title formatting (if we pass a title string)
    if title_text is not None:
        title = slide.shapes.title
        title.text = title_text
        tf = title.text_frame.paragraphs[0]
        tf.font.name = "Arial"
        tf.font.size = Pt(36)
        tf.font.bold = True
        tf.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)  # dark navy
        tf.alignment = PP_ALIGN.CENTER

# -------------------------------------------------------------------
# Output location (repo root / docs)
# -------------------------------------------------------------------
OUTPUT_DIR = Path(__file__).resolve().parents[1] / "docs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
output_path = OUTPUT_DIR / "DriveLegal_AI_7Slide_Presentation.pptx"

prs = Presentation()

# -------------------------------------------------------------------
# Slide builders
# -------------------------------------------------------------------
def add_title_slide(title, subtitle="Road Safety Hackathon 2026 – CoERS, IIT Madras × MoRTH"):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    apply_theme(slide, title)
    # subtitle placeholder (layout 0 has a second placeholder)
    slide.placeholders[1].text = subtitle
    # subtitle styling: lighter font
    sub_tf = slide.placeholders[1].text_frame.paragraphs[0]
    sub_tf.font.name = "Arial"
    sub_tf.font.size = Pt(24)
    sub_tf.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    sub_tf.alignment = PP_ALIGN.CENTER

def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    apply_theme(slide)
    slide.shapes.title.text = title
    # Title styling (reuse same style as apply_theme)
    title_tf = slide.shapes.title.text_frame.paragraphs[0]
    title_tf.font.name = "Arial"
    title_tf.font.size = Pt(30)
    title_tf.font.bold = True
    title_tf.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.font.name = "Arial"
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # dark gray for readability
        p.level = 0

def add_image_slide(title, image_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title & Content layout
    apply_theme(slide)
    slide.shapes.title.text = title
    # Title styling
    title_tf = slide.shapes.title.text_frame.paragraphs[0]
    title_tf.font.name = "Arial"
    title_tf.font.size = Pt(30)
    title_tf.font.bold = True
    title_tf.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    slide.shapes.add_picture(image_path, left, top, width=width)
    if caption:
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.name = "Arial"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
        p.alignment = PP_ALIGN.CENTER

# -------------------------------------------------------------------
# Build the 7‑slide deck
# -------------------------------------------------------------------
# 1. Welcome
add_title_slide("🚦 Welcome to DriveLegal AI", "AI‑powered traffic‑law assistant")

# 2. Problem Statement
add_bullet_slide("The Problem", [
    "Millions of Indian drivers receive over‑charged traffic challans every year.",
    "Legal fine amounts differ across 18+ states – drivers lack clear guidance.",
    "Disputing fines requires a lawyer – costly and time‑consuming.",
    "Language barriers leave rural users unsupported."
])

# 3. Our Solution
add_bullet_slide("Our Solution", [
    "Upload a challan photo – AI instantly validates the fine against the MV Act 2019.",
    "One‑click generation of a printer‑ready dispute letter with legal citations.",
    "Multilingual AI chat (10+ Indian languages) for any traffic‑law query.",
    "Integrated calculators: fine, BAC, penalty points, document expiry, speed limits."
])

# 4. Demo – Over‑charge Detection (image)
add_image_slide(
    "Demo: Over‑charge Detection",
    "C:/Users/rohit/.gemini/antigravity/brain/ffefbd62-1568-44d6-ad8b-3d5e37b1379f/overcharge_demo_1779964929978.png",
    "AI flags ₹5,000 fine as OVERCHARGED – correct fine is ₹1,000."
)

# 5. Impact & Benefits
add_bullet_slide("Impact", [
    "💰 Saves drivers up to ₹50,000 per year by preventing over‑charges.",
    "🛡️ Empowers 30 crore Indian vehicle owners with legal clarity.",
    "🚫 Reduces corruption – officers cannot arbitrarily inflate fines.",
    "🤝 Aligns with MoRTH Vision Zero 2030 and Digital India initiatives."
])

# 6. Roadmap & Future Extensions
add_bullet_slide("Roadmap", [
    "WhatsApp bot for on‑the‑go access.",
    "GPS‑based automatic state detection.",
    "Integration with UMANG & DigiLocker for document verification.",
    "AI‑driven outcome predictor for court cases."
])

# 7. Thank You
add_title_slide(
    "Thank You!",
    "Visit the live demo: https://rohitheevlsi-drivelegal-ai.onrender.com\nGitHub: https://github.com/rohitheevlsi/rohitheevlsi-drivelegal-ai"
)

# Save the presentation
prs.save(output_path)
print(f"7-slide presentation generated at {output_path}")
