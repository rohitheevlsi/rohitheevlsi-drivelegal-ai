# generate_presentation.py
"""Generate a powerful, hackathon‑winning PowerPoint deck for DriveLegal AI.

Run this script (requires python‑pptx) to create `DriveLegal_AI_Presentation.pptx` in the `docs/` folder.

    pip install python-pptx
    python scripts/generate_presentation.py
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Output path (repo root/docs)
OUTPUT_DIR = Path(__file__).resolve().parents[1] / "docs"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
output_path = OUTPUT_DIR / "DriveLegal_AI_Presentation.pptx"

prs = Presentation()

# ----- Helper -----
def add_title_slide(title, subtitle="Road Safety Hackathon 2026 – CoERS, IIT Madras × MoRTH"):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_tf = slide.shapes.title.text_frame
    title_tf.text = title
    subtitle_tf = slide.placeholders[1].text_frame
    subtitle_tf.text = subtitle

def add_content_slide(title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    for bp in bullet_points:
        p = tf.add_paragraph()
        p.text = bp
        p.font.size = Pt(24)
        p.level = 0

# ----- Deck -----
add_title_slide("🚦 DriveLegal AI", "AI‑powered traffic‑law assistant")

add_content_slide("Problem Landscape", [
    "Millions of over‑charged traffic challans every year",
    "Drivers lack legal knowledge of the Motor Vehicles Act 2019",
    "State‑wise fine variations create confusion",
    "Manual dispute letters cost time & money"
])

add_content_slide("Our Solution", [
    "Upload a challan photo – AI validates legality instantly",
    "One‑click generation of a printer‑ready dispute letter",
    "Multilingual AI chat for any traffic‑law query (10+ Indian languages)",
    "State‑wise fine calculator, BAC estimator, penalty‑points tracker"
])

add_content_slide("Key Features (10 Tools)", [
    "AI Chat", "Challan Validator", "Dispute Letter Generator", "Fine Calculator",
    "BAC Calculator", "Penalty‑Points Tracker", "Document Expiry Checker",
    "State Rules Comparator", "Rights Advisor", "Speed‑Limits Guide"
])

add_content_slide("Tech Stack", [
    "Frontend: Streamlit (Python) – rapid UI", "AI backend: Google Gemini (Gemini‑3.5‑flash)",
    "Data: Embedded MV Act 2019 tables for 18 states", "No external DB – fully offline for calculators",
    "Deployment: Render (free tier) – zero‑config CI/CD"
])

add_content_slide("Demo Walkthrough (GIF)", [
    "1️⃣ Upload challan – AI flags over‑charge (₹5,000 → ₹1,000)",
    "2️⃣ Auto‑generate dispute letter with legal citations",
    "3️⃣ Use multilingual chat to ask any traffic‑law question",
    "4️⃣ Visualize state‑wise fine comparison"
])

add_content_slide("Real‑World Impact", [
    "📊 Saves ₹10‑50 k per driver by preventing over‑charges", "🛡️ Empowers 30 crore Indian drivers with legal clarity",
    "🚫 Reduces corruption – officers cannot arbitrarily inflate fines",
    "🤝 Aligns with MoRTH’s Vision Zero 2030 goals"
])

add_content_slide("Roadmap & Future Extensions", [
    "WhatsApp bot – on‑the‑go access", "GPS‑based state auto‑detect", "Integration with UMANG & DigiLocker",
    "AI‑driven outcome predictor for court cases"
])

add_content_slide("Team & Credits", [
    "👤 Rohit Heevlsi – Lead Engineer & Demo creator", "🤝 Co‑authors: IIT Madras, MoRTH advisors",
    "🔧 Open‑source on GitHub – contributions welcome"
])

add_content_slide("Thank You!", ["Visit the live demo: https://rohitheevlsi-drivelegal-ai.onrender.com", "GitHub repo: https://github.com/rohitheevlsi/rohitheevlsi-drivelegal-ai"])

# Save the PPTX
prs.save(output_path)
print(f"Presentation generated at {output_path}")
